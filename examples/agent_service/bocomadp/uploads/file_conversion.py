# -*- coding: utf-8 -*-
"""上传文件格式转换（移植自 deer-flow utils/file_conversion.py）。

统一把可文本化的文件转换为 **Markdown (.md)** 文本。与历史版本不同，
本模块不再直接落盘——转换在 **host 侧**（第三方库）完成，返回 markdown
文本字符串，由上层（routers/uploads.py）经 ``workspace.get_backend()``
写入沙箱；该设计使上传逻辑沙箱感知（双 PVC / 共享 PVC 下 session 隔离）。

支持：
- 文本/代码类（txt/md/csv/json/xml/log/各类源码） -> 复制为同名 .md
- PDF  -> .md (pdfplumber)
- Word (.docx/.doc) -> .md (python-docx)
- PPT  (.pptx/.ppt) -> .md (python-pptx)
- Excel (.xlsx/.xls) -> .md 表格 (openpyxl/pandas)
- HTML -> .md (html2text)
不支持：图片、压缩包、二进制等 -> 由调用方按需拒绝。
"""
from __future__ import annotations

import io

from .manager import UploadError


# 扩展名 -> (类别) 映射
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".xml",
    ".log", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".sh", ".bash",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".go", ".c", ".cpp", ".h",
    ".hpp", ".cs", ".rb", ".php", ".rs", ".kt", ".swift", ".sql", ".r", ".scala",
    ".pl", ".lua", ".vim", ".dockerfile", ".gitignore", ".env",
}
_PDF_EXTS = {".pdf"}
_WORD_EXTS = {".docx", ".doc"}
_PPT_EXTS = {".pptx", ".ppt"}
_EXCEL_EXTS = {".xlsx", ".xls", ".xlsm"}
_HTML_EXTS = {".html", ".htm"}


class UnsupportedFileType(UploadError):
    """文件类型不在支持范围内。"""


def is_supported_format(filename: str) -> bool:
    """判断文件名是否可转换为 .md。"""
    ext = f".{_split_ext(filename)}"  # 补点后与下方带点扩展名集合比较
    return any(
        ext in group
        for group in (
            _TEXT_EXTS, _PDF_EXTS, _WORD_EXTS, _PPT_EXTS, _EXCEL_EXTS, _HTML_EXTS,
        )
    )


def _split_ext(filename: str) -> str:
    name = (filename or "").strip().replace("\\", "/").rsplit("/", 1)[-1]
    return name.lower().rsplit(".", 1)[-1] if "." in name else ""


def convert_file_bytes(
    filename: str,
    data: bytes,
    content_type: str | None = None,
) -> tuple[str, str]:
    """把上传文件的字节内容转换为 markdown 文本（host 侧执行）。

    Args:
        filename: 客户端原始文件名（含扩展名）。
        data: 已读取的文件字节。
        content_type: 可选 MIME 类型（当前仅作候选判定，未强制）。

    Returns:
        ``(format, markdown_text)``；无法转换时抛
        `UnsupportedFileType` 或具体转换错误（调用方捕获后不阻断上传）。
    """
    ext = f".{_split_ext(filename)}"
    if ext in _TEXT_EXTS:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="ignore")
        return "text", text

    if ext in _PDF_EXTS:
        return "pdf", _convert_pdf(data)
    if ext in _WORD_EXTS:
        return "word", _convert_docx(data)
    if ext in _PPT_EXTS:
        return "ppt", _convert_pptx(data)
    if ext in _EXCEL_EXTS:
        return "excel", _convert_excel(data)
    if ext in _HTML_EXTS:
        return "html", _convert_html(data)

    raise UnsupportedFileType(f"unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# 具体格式转换实现（按需 import 第三方库，缺失则标记不支持）
# 统一签名：(src: bytes) -> md_text: str
# ---------------------------------------------------------------------------
def _convert_pdf(src: bytes) -> str:
    try:
        import pdfplumber  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("pdfplumber not installed") from e
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(src)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            parts.append(f"## 第 {i} 页\n\n{text}")
    return "\n\n".join(parts)


def _convert_docx(src: bytes) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-docx not installed") from e
    doc = Document(io.BytesIO(src))
    lines: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "") if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            level = "".join(filter(str.isdigit, style)) or "1"
            lines.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def _convert_pptx(src: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-pptx not installed") from e
    prs = Presentation(io.BytesIO(src))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        lines.append(t)
    return "\n\n".join(lines)


def _convert_excel(src: bytes) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("openpyxl not installed") from e
    wb = openpyxl.load_workbook(io.BytesIO(src), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"## Sheet: {ws.title}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else "" for c in rows[0]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def _convert_html(src: bytes) -> str:
    try:
        import html2text  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("html2text not installed") from e
    h = html2text.HTML2Text()
    h.body_width = 0  # 不自动换行
    raw = src.decode("utf-8", errors="ignore")
    return h.handle(raw)
